import pptx
from pptx import Presentation
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.util import Cm, Pt, Inches
from pptx.dml.color import RGBColor
from lxml import etree
import glob,os

def settext(text_frame,text,fontname='Ariel',size=16,bold=False,color=[0,0,0],bulleted=False,clear=False):
    if clear:
        text_frame.clear()  # not necessary for newly-created shape
    p = text_frame.paragraphs[0]
    # p.font.size = pptx.util.Pt(size)
    # p.font.bold = bold
    run = p.add_run()
    run.text = text
    font = run.font
    font.name = fontname
    font.size = Pt(size)
    font.bold = bold
    font.color.rgb = RGBColor(*color)
    if not bulleted:
        p._pPr.insert(0, etree.Element("{http://schemas.openxmlformats.org/drawingml/2006/main}buNone"))
def addadvrslide(prs,title,text,imagefile,proprietary=True):
    title_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(title_slide_layout)
    subtitle = slide.placeholders[1]
    #add default elements
    settext(slide.shapes.title.text_frame,title+'\n',fontname='Times',size=32,bold=True,color=[0x33,0x33,0x99])
    separator = slide.shapes.add_picture('c:/py/work/pypowerpoint/rainbow.png', left=Inches(0), top=Inches(1.2), width=Inches(10))
    logo = slide.shapes.add_picture('c:/py/work/pypowerpoint/advrlogo.png', left=Inches(0.16), top=Inches(7), width=Inches(1.43))
    box = slide.shapes.add_textbox(left=Inches(8.23), top=Inches(7.27), width=Inches(1.54), height=Inches(0.24))
    if proprietary:
        settext(box.text_frame,'ADVR Inc Proprietary',size=10)
    #add content
    image = slide.shapes.add_picture(imagefile,left=Inches(0),top=Inches(2))
    image.left = (prs.slide_width - image.width) // 2
    if text:
        settext(subtitle.text_frame,text)
    else:
        sp = slide.shapes[1].element
        sp.getparent().remove(sp)
    # send image to back https://github.com/scanny/python-pptx/issues/195
    slide.shapes._spTree.remove(image._element)
    slide.shapes._spTree.insert(2, image._element)
def addtestslide(prs):
    title_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    # title.text = "Hello World!"
    # subtitle.text = "python-pptx was here!"
    settext(title.text_frame,'Title'+'\n',fontname='Times',size=32,bold=True,color=[0x33,0x33,0x99])
    separator = slide.shapes.add_picture('rainbow.png', left=Inches(0), top=Inches(1.2), width=Inches(10))
    logo = slide.shapes.add_picture('advrlogo.png', left=Inches(0.16), top=Inches(7), width=Inches(1.43))
    box = slide.shapes.add_textbox(left=Inches(8.23), top=Inches(7.27), width=Inches(1.54), height=Inches(0.24))
    settext(box.text_frame,'ADVR Inc Proprietary',size=10)
    image = slide.shapes.add_picture('h11.png',left=Inches(0),top=Inches(2))
    image.left = (prs.slide_width - image.width) // 2
    settext(subtitle.text_frame,'Interleaved period = 25.7 µm\nInterleaved period = 35.7 µm',bulleted=True)
    slide.shapes._spTree.remove(image._element)
    slide.shapes._spTree.insert(2, image._element)
def analyze_ppt(input, output): # https://pbpython.com/creating-powerpoint.html
    """ Take the input file and analyze the structure.
    The output file contains marked up information to make it easier
    for generating future powerpoint templates.
    """
    prs = Presentation(input)
    # Each powerpoint file has multiple layouts
    # Loop through them all and  see where the various elements are
    for index, _ in enumerate(prs.slide_layouts):
        slide = prs.slides.add_slide(prs.slide_layouts[index])
        # Not every slide has to have a title
        try:
            title = slide.shapes.title
            title.text = 'Title for Layout {}'.format(index)
        except AttributeError:
            print("No Title for Layout {}".format(index))
        # Go through all the placeholders and identify them by index and type
        for shape in slide.placeholders:
            if shape.is_placeholder:
                phf = shape.placeholder_format
                # Do not overwrite the title which is just a special placeholder
                try:
                    if 'Title' not in shape.text:
                        shape.text = 'Placeholder index:{} type:{}'.format(phf.idx, shape.name)
                except AttributeError:
                    print("{} has no text attribute".format(phf.type))
                print('{} {}'.format(phf.idx, shape.name))
    prs.save(output)

if __name__ == '__main__':
    def test():
        prs = Presentation()
        # prs = Presentation('advr.pptx')
        title_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = "Hello, World!"
        subtitle.text = "python-pptx was here!"
        prs.save('test.pptx')
    def newtest(): # default slide size = 10" x 7.5"
        os.system('taskkill /F /FI "WINDOWTITLE eq newtest.pptx*" ')
        prs = Presentation()
        addtestslide(prs)
        prs.save('newtest.pptx')
        os.startfile('C:\\py\\work\\pypowerpoint\\newtest.pptx')
    # analyze_ppt('advr.pptx','test.pptx')
    # test()
    newtest()
